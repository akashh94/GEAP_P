import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette Constants ──
BG_COLOR = RGBColor(12, 10, 18)        # Very dark navy-purple background (#0C0A12)
CARD_BG = RGBColor(27, 23, 40)         # Glass/Card dark purple (#1B1728)
ACCENT_PURPLE = RGBColor(147, 51, 234) # Electric Purple (#9333EA)
ACCENT_GOLD = RGBColor(197, 160, 89)   # Wealth management gold (#C5A059)
TEXT_WHITE = RGBColor(248, 250, 252)   # Title/Bright text (#F8FAFC)
TEXT_MUTED = RGBColor(148, 163, 184)   # Secondary label text (#94A3B8)
TEXT_GREEN = RGBColor(52, 211, 153)    # Success / Positive gains (#34D399)
TEXT_AMBER = RGBColor(245, 158, 11)    # Risk / Warning (#F59E0B)

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_header(slide, title_text, category_text="GEAP E*TRADE POC"):
    # Category / Breadcrumb
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = "Outfit"
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_GOLD
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "Outfit"
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE

def create_card(slide, left, top, width, height, title="", border_color=None):
    # Draw background rectangle shape representing a card/widget
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
    else:
        card.line.fill.background() # borderless
        
    if title:
        # Text box for title
        tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Outfit"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GOLD
        
    return card

def main():
    prs = Presentation()
    
    # Set slide size to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # ═══════════════════════════════════════════════════════════════
    # SLIDE 1: Title Slide (Dark Cover)
    # ═══════════════════════════════════════════════════════════════
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)
    
    # Title & Subtitle in single textbox
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "GEMINI ENTERPRISE AGENT PLATFORM"
    p1.font.name = "Outfit"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = "E*TRADE Wealth Integration Proof-of-Concept"
    p2.font.name = "Outfit"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_GOLD
    p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = "Secure Autonomous Multi-Agent Workflows in Wealth Management"
    p3.font.name = "Outfit"
    p3.font.size = Pt(16)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(20)
    
    p4 = tf.add_paragraph()
    p4.text = "Executive Summary & System Architecture Deck"
    p4.font.name = "Outfit"
    p4.font.size = Pt(12)
    p4.font.bold = True
    p4.font.color.rgb = ACCENT_PURPLE
    p4.space_before = Pt(30)
    
    # Decorative Accent Line
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(5.8), Inches(4.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_PURPLE
    line.line.fill.background()

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 2: Exec Summary / Business Challenge & Opportunity
    # ═══════════════════════════════════════════════════════════════
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "Executive Summary & Core Objective")
    
    # Left Card: The Challenge
    create_card(slide2, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "THE BUSINESS CHALLENGE")
    tb_chall = slide2.shapes.add_textbox(Inches(0.95), Inches(2.3), Inches(5.3), Inches(4.0))
    tf_chall = tb_chall.text_frame
    tf_chall.word_wrap = True
    
    p = tf_chall.paragraphs[0]
    p.text = "Self-directed brokerage accounts represent high compliance, privacy, and cybersecurity risks for autonomous AI deployments. Key issues include:"
    p.font.name = "Outfit"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(8)
    
    bullets = [
        "PII Exposure: Leaking social security and account credentials to public LLMs.",
        "Fiduciary Liability: Providing automated advisory suggestions without mandatory compliance disclaimers.",
        "Unregulated Transfers: Risks associated with automated cash movement or wire fraud.",
        "UTMA Minor Account Lockouts: Compliance gates protecting beneficiary wealth."
    ]
    for b in bullets:
        p_b = tf_chall.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = "Outfit"
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = TEXT_MUTED
        p_b.space_before = Pt(6)
        
    # Right Card: The Solution
    create_card(slide2, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8), "THE GEAP POC SOLUTION")
    tb_sol = slide2.shapes.add_textbox(Inches(7.05), Inches(2.3), Inches(5.3), Inches(4.0))
    tf_sol = tb_sol.text_frame
    tf_sol.word_wrap = True
    
    p = tf_sol.paragraphs[0]
    p.text = "A framework-free proof of concept demonstrating Morgan Stanley Wealth integration boundaries. Key highlights:"
    p.font.name = "Outfit"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(8)
    
    sol_bullets = [
        "Security-First Middleware: Scans inputs client-side, redacting sensitive assets before LLM delivery.",
        "Dynamic Auto-Routing: Multi-agent fleet (7 nodes) handles tasks under unique instructions and whitelists.",
        "Fiduciary Guardrails: Automatically appends SIPC alerts and blocks wire movements to enforce manual callback loops.",
        "No-Framework Performance: Ultra-fast load times with vanilla JS and local server-side minification."
    ]
    for sb in sol_bullets:
        p_sb = tf_sol.add_paragraph()
        p_sb.text = "✦ " + sb
        p_sb.font.name = "Outfit"
        p_sb.font.size = Pt(12)
        p_sb.font.color.rgb = TEXT_MUTED
        p_sb.space_before = Pt(6)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 3: Secure Architecture Pipeline
    # ═══════════════════════════════════════════════════════════════
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "Secure Prompt Processing Pipeline")
    
    # 5 Process Steps horizontally
    steps = [
        ("1. Input & Scrub", "PII Shield", "Detects and redacts SSNs, phone numbers, emails, and accounts before LLM transmission."),
        ("2. Intent Routing", "Auto-Router", "Scores queries using keyword weights & active page bias to assign the best specialist agent."),
        ("3. Context Injection", "Context Engine", "Appends real-time routes, balances, whitelist sources, and read-only restrictions."),
        ("4. Stream Engine", "Gemini 2.0 Flash", "Processes requests and handles recursive function-calling loops locally."),
        ("5. Output Guard", "Compliance Gate", "Appends investment disclaimers, blocks outbound wires, and stages human desk fallbacks.")
    ]
    
    card_width = Inches(2.2)
    card_height = Inches(4.4)
    spacing = Inches(0.2)
    start_left = Inches(0.8)
    
    for i, (num_title, label, desc) in enumerate(steps):
        left = start_left + i * (card_width + spacing)
        # Highlight card 5 for security focus
        border = ACCENT_PURPLE if i == 4 or i == 0 else None
        create_card(slide3, left, Inches(2.0), card_width, card_height, num_title, border)
        
        # Details inside card
        tb = slide3.shapes.add_textbox(left + Inches(0.1), Inches(2.5), card_width - Inches(0.2), Inches(3.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = label
        p.font.name = "Outfit"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(10)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Outfit"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(8)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 4: Core Engines (Context & Router)
    # ═══════════════════════════════════════════════════════════════
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "Unified Context Engine & Auto-Router")
    
    # Left Card: Context Engine
    create_card(slide4, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "UNIFIED CONTEXT ENGINE (context.js)")
    tb_ctx = slide4.shapes.add_textbox(Inches(0.95), Inches(2.3), Inches(5.3), Inches(4.0))
    tf_ctx = tb_ctx.text_frame
    tf_ctx.word_wrap = True
    
    p = tf_ctx.paragraphs[0]
    p.text = "Ensures zero context drift across multiple agent instances. Injects a structured state package into the system prompt:"
    p.font.name = "Outfit"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(10)
    
    ctx_bullets = [
        "Selected Account context (net values, settling cash, account type).",
        "Active platform details (route pathname, CSS theme).",
        "User permissions and whitelisted sources (e.g. Yahoo Finance, SEC Edgar).",
        "Read-Only status to block write/navigation tools in restrictive modes."
    ]
    for cb in ctx_bullets:
        p_cb = tf_ctx.add_paragraph()
        p_cb.text = "✔ " + cb
        p_cb.font.name = "Outfit"
        p_cb.font.size = Pt(12)
        p_cb.font.color.rgb = TEXT_MUTED
        p_cb.space_before = Pt(6)
        
    # Right Card: Auto Router
    create_card(slide4, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8), "INTENT AUTO-ROUTER (agents.js)")
    tb_rt = slide4.shapes.add_textbox(Inches(7.05), Inches(2.3), Inches(5.3), Inches(4.0))
    tf_rt = tb_rt.text_frame
    tf_rt.word_wrap = True
    
    p = tf_rt.paragraphs[0]
    p.text = "Routes queries instantly without requesting user confirmation. It processes routing weight through a two-fold check:"
    p.font.name = "Outfit"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(10)
    
    rt_bullets = [
        "Keyword Matches: Matches queries with thematic keywords (e.g. 'tax harvest' -> Analyst, 'refi' -> Mortgage).",
        "Active Page Route Bias: Shifts auto-routing weights dynamically based on what page the user is viewing.",
        "Dynamic Fleet Management: Deploys register/deregister APIs allowing admins to spawn/retract agents on-the-fly."
    ]
    for rb in rt_bullets:
        p_rb = tf_rt.add_paragraph()
        p_rb.text = "✦ " + rb
        p_rb.font.name = "Outfit"
        p_rb.font.size = Pt(12)
        p_rb.font.color.rgb = TEXT_MUTED
        p_rb.space_before = Pt(6)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 5: Specialized AI Fleet (Agents Grid)
    # ═══════════════════════════════════════════════════════════════
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "The Specialized Multi-Agent Fleet")
    
    # 6 cards in a 2x3 grid
    agents_data = [
        ("Portfolio Analyst 📈", "Allocation & Risk", "Reviews weights, sector concentration, and suggests rebalancing and tax-loss harvesting candidates."),
        ("Trade Assistant ⚡", "Execution & Quotes", "Provides quotes, explains limit/stop orders, size risk, and validates buying power."),
        ("Market Research 🔬", "Market Intelligence", "Reviews macro data, trends, stock comparisons, and fundamentals (P/E, EPS)."),
        ("Customer Support 💬", "Navigation & FAQ", "Answers platform support, transfer methods, documents, and fee structures."),
        ("Super Agent 🔬", "MS Research & Alts", "Integrates Morgan Stanley Research and alts (crypto/ETFs, gold, private equity)."),
        ("Mortgage Agent 🏠", "Mortgage & HELOC", "Directs Jumbo/Rocket paths and extracts customer location, timelines, and HELOC targets.")
    ]
    
    w_card = Inches(3.6)
    h_card = Inches(2.2)
    s_x = Inches(0.4)
    s_y = Inches(0.3)
    
    for i, (name, role, desc) in enumerate(agents_data):
        row = i // 3
        col = i % 3
        
        left = Inches(0.8) + col * (w_card + s_x)
        top = Inches(1.8) + row * (h_card + s_y)
        
        create_card(slide5, left, top, w_card, h_card, name)
        
        tb = slide5.shapes.add_textbox(left + Inches(0.15), top + Inches(0.5), w_card - Inches(0.3), Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = role.upper()
        p.font.name = "Outfit"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = ACCENT_PURPLE
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Outfit"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(4)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 6: Privacy Shield & PII Scrubber
    # ═══════════════════════════════════════════════════════════════
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "Privacy Shield & PII Scrubber")
    
    # Left Card: How PII Scrubber works
    create_card(slide6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "CLIENT-SIDE DATA REDACTION")
    tb_pii = slide6.shapes.add_textbox(Inches(0.95), Inches(2.3), Inches(5.3), Inches(4.0))
    tf_pii = tb_pii.text_frame
    tf_pii.word_wrap = True
    
    p = tf_pii.paragraphs[0]
    p.text = "Before transmitting prompt text to external Google Gemini servers, the GEAP local middleware scans and scrubs sensitive data matching:"
    p.font.name = "Outfit"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(10)
    
    pii_types = [
        "Social Security Numbers (redacted to [REDACTED SSN])",
        "Brokerage Account Numbers (redacted to [REDACTED ACCOUNT])",
        "Email Addresses (redacted to [REDACTED EMAIL])",
        "Phone Numbers (redacted to [REDACTED PHONE])"
    ]
    for pt in pii_types:
        p_pt = tf_pii.add_paragraph()
        p_pt.text = "• " + pt
        p_pt.font.name = "Outfit"
        p_pt.font.size = Pt(12)
        p_pt.font.color.rgb = TEXT_MUTED
        p_pt.space_before = Pt(6)
        
    p_last = tf_pii.add_paragraph()
    p_last.text = "🛡️ Prompts reach the cloud clean. This ensures strict compliance with Morgan Stanley data governance guidelines."
    p_last.font.name = "Outfit"
    p_last.font.size = Pt(12)
    p_last.font.bold = True
    p_last.font.color.rgb = TEXT_GREEN
    p_last.space_before = Pt(20)

    # Right Card: Demo Example
    create_card(slide6, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8), "DEMONSTRATION RUN")
    tb_demo = slide6.shapes.add_textbox(Inches(7.05), Inches(2.3), Inches(5.3), Inches(4.0))
    tf_demo = tb_demo.text_frame
    tf_demo.word_wrap = True
    
    p = tf_demo.paragraphs[0]
    p.text = "User Input:"
    p.font.name = "Outfit"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    p.space_after = Pt(2)
    
    p_in = tf_demo.add_paragraph()
    p_in.text = "\"Test PII scrubber: my SSN is 111-22-3333 and my account is 1234-5678\""
    p_in.font.name = "Consolas"
    p_in.font.size = Pt(11)
    p_in.font.color.rgb = TEXT_WHITE
    p_in.space_before = Pt(2)
    p_in.space_after = Pt(12)
    
    p_trans = tf_demo.add_paragraph()
    p_trans.text = "Scrubbed Prompt Transmitted to LLM:"
    p_trans.font.name = "Outfit"
    p_trans.font.size = Pt(11)
    p_trans.font.bold = True
    p_trans.font.color.rgb = ACCENT_PURPLE
    
    p_out = tf_demo.add_paragraph()
    p_out.text = "\"Test PII scrubber: my SSN is [REDACTED SSN] and my account is [REDACTED ACCOUNT]\""
    p_out.font.name = "Consolas"
    p_out.font.size = Pt(11)
    p_out.font.color.rgb = TEXT_GREEN
    p_out.space_before = Pt(2)
    p_out.space_after = Pt(12)
    
    p_desc = tf_demo.add_paragraph()
    p_desc.text = "Visual Alert: The interface immediately renders a security banner informing the user that their data has been scrubbed for privacy before sending."
    p_desc.font.name = "Outfit"
    p_desc.font.size = Pt(12)
    p_desc.font.color.rgb = TEXT_MUTED

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 7: Compliance Gates & Human Escalation
    # ═══════════════════════════════════════════════════════════════
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "Compliance Gates & Human Handoffs")
    
    # Left Card: Fiduciary & UTMA
    create_card(slide7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "FIDUCIARY DISCLOSURES & UTMA PROTECTION")
    tb_fid = slide7.shapes.add_textbox(Inches(0.95), Inches(2.3), Inches(5.3), Inches(4.0))
    tf_fid = tb_fid.text_frame
    tf_fid.word_wrap = True
    
    p = tf_fid.paragraphs[0]
    p.text = "Fiduciary Advice Disclaimers:"
    p.font.name = "Outfit"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    p_fid_desc = tf_fid.add_paragraph()
    p_fid_desc.text = "When investment keywords (e.g. portfolio, buy, rebalance, fund) are detected, the system appends mandatory disclaimers detailing that advice is educational and not personalized."
    p_fid_desc.font.name = "Outfit"
    p_fid_desc.font.size = Pt(11)
    p_fid_desc.font.color.rgb = TEXT_MUTED
    p_fid_desc.space_before = Pt(4)
    p_fid_desc.space_after = Pt(12)
    
    p_utma = tf_fid.add_paragraph()
    p_utma.text = "UTMA Minor Account Lockouts:"
    p_utma.font.name = "Outfit"
    p_utma.font.size = Pt(13)
    p_utma.font.bold = True
    p_utma.font.color.rgb = TEXT_WHITE
    
    p_utma_desc = tf_utma_desc = tf_fid.add_paragraph()
    p_utma_desc.text = "If a custodial UTMA account is selected, outbound electronic cash transfers are blocked. System raises UTMA compliance gates reminding that custodial funds belong irrevocably to the minor."
    p_utma_desc.font.name = "Outfit"
    p_utma_desc.font.size = Pt(11)
    p_utma_desc.font.color.rgb = TEXT_MUTED
    p_utma_desc.space_before = Pt(4)

    # Right Card: Human Handoff
    create_card(slide7, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8), "HIGH-RISK ACTION HANDOFF")
    tb_esc = slide7.shapes.add_textbox(Inches(7.05), Inches(2.3), Inches(5.3), Inches(4.0))
    tf_esc = tb_esc.text_frame
    tf_esc.word_wrap = True
    
    p = tf_esc.paragraphs[0]
    p.text = "Wire & Withdrawal Gates:"
    p.font.name = "Outfit"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    p_esc_desc = tf_esc.add_paragraph()
    p_esc_desc.text = "Outbound wires or transfers above $10,000 are classified as high-risk. Autonomous execution is blocked. Instead, the system stages a transaction draft and initiates handoff:"
    p_esc_desc.font.name = "Outfit"
    p_esc_desc.font.size = Pt(11)
    p_esc_desc.font.color.rgb = TEXT_MUTED
    p_esc_desc.space_before = Pt(4)
    p_esc_desc.space_after = Pt(10)
    
    esc_steps = [
        "Intent captured and sent to the core advisor dashboard.",
        "Generates support callback tickets (e.g. MS-WIRE-84920).",
        "Directs the user to secure telephone advisor desks."
    ]
    for es in esc_steps:
        p_es = tf_esc.add_paragraph()
        p_es.text = "➔ " + es
        p_es.font.name = "Outfit"
        p_es.font.size = Pt(11)
        p_es.font.color.rgb = TEXT_MUTED
        p_es.space_before = Pt(4)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 8: Admin Tools (Studio & Fabric)
    # ═══════════════════════════════════════════════════════════════
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_header(slide8, "Admin Diagnostics & Monitoring Suite")
    
    # Left Card: Agent Studio
    create_card(slide8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "AGENT STUDIO (CONFIGURATION)")
    tb_st = slide8.shapes.add_textbox(Inches(0.95), Inches(2.3), Inches(5.3), Inches(4.0))
    tf_st = tb_st.text_frame
    tf_st.word_wrap = True
    
    p = tf_st.paragraphs[0]
    p.text = "A dashboard allowing admins to design, configure, and inspect individual agents:"
    p.font.name = "Outfit"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(10)
    
    st_bullets = [
        "Visual Node Flow: Displays logic pathways from User Input -> Router -> Agent Fleet.",
        "Model Customization: Adjusts model templates and temperature parameters.",
        "Equipped Tool Whitelisting: Selects specific JS tools allowed for each agent.",
        "Studio Test Console: Offers isolated text testing with context print-outs."
    ]
    for sb in st_bullets:
        p_sb = tf_st.add_paragraph()
        p_sb.text = "• " + sb
        p_sb.font.name = "Outfit"
        p_sb.font.size = Pt(12)
        p_sb.font.color.rgb = TEXT_MUTED
        p_sb.space_before = Pt(6)
        
    # Right Card: Agent Fabric
    create_card(slide8, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8), "AGENT FABRIC (GOVERNANCE)")
    tb_fb = slide8.shapes.add_textbox(Inches(7.05), Inches(2.3), Inches(5.3), Inches(4.0))
    tf_fb = tb_fb.text_frame
    tf_fb.word_wrap = True
    
    p = tf_fb.paragraphs[0]
    p.text = "Widescreen analytics console showing real-time agent metrics and controls:"
    p.font.name = "Outfit"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(10)
    
    fb_bullets = [
        "Performance Monitoring: Tracks total requests, response latency, token consumption, and model success rates.",
        "Real-Time Activity Log: Visualizes live routed queries and error status flags.",
        "Active Governance Panel: Allows toggling safety guardrails, PII scrubbers, and disclaimer rules.",
        "Router Mapping: Inspects network loads across agents in the fleet."
    ]
    for fb in fb_bullets:
        p_fb = tf_fb.add_paragraph()
        p_fb.text = "• " + fb
        p_fb.font.name = "Outfit"
        p_fb.font.size = Pt(12)
        p_fb.font.color.rgb = TEXT_MUTED
        p_fb.space_before = Pt(6)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 9: Tech Stack & System Verification
    # ═══════════════════════════════════════════════════════════════
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_header(slide9, "Tech Stack & System Verification")
    
    # Left Card: Technology Stack
    create_card(slide9, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "LIGHTWEIGHT STACK SUMMARY")
    tb_stack = slide9.shapes.add_textbox(Inches(0.95), Inches(2.3), Inches(5.3), Inches(4.0))
    tf_stack = tb_stack.text_frame
    tf_stack.word_wrap = True
    
    p = tf_stack.paragraphs[0]
    p.text = "The application avoids bulky frameworks to minimize dependency risk and simplify audits:"
    p.font.name = "Outfit"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(10)
    
    st_items = [
        "Backend: Node.js, Express, OAuth 1.0a handshake, custom server.log routing.",
        "Frontend: HTML5, CSS3, Vanilla JS SPA framework, Canvas charting.",
        "LLM SDK: @google/genai fetched dynamically via esm.run, streaming.",
        "News Feed Proxy: Node parsing Google News XML feed into JSON."
    ]
    for s_i in st_items:
        p_si = tf_stack.add_paragraph()
        p_si.text = "✦ " + s_i
        p_si.font.name = "Outfit"
        p_si.font.size = Pt(12)
        p_si.font.color.rgb = TEXT_MUTED
        p_si.space_before = Pt(6)

    # Right Card: System Verification
    create_card(slide9, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8), "TEST COVERAGE")
    tb_test = slide9.shapes.add_textbox(Inches(7.05), Inches(2.3), Inches(5.3), Inches(4.0))
    tf_test = tb_test.text_frame
    tf_test.word_wrap = True
    
    p = tf_test.paragraphs[0]
    p.text = "A full automated verification test harness ensures deployment stability:"
    p.font.name = "Outfit"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(10)
    
    test_bullets = [
        "Unit Tests: Validates auto-routing keyword scores, transaction calculations, and context engine payloads.",
        "Integration Tests: Validates session storage updates, Google News proxy parsing, and Express endpoints.",
        "E2E Browser Tests: Puppeteer scripts launch a headless browser, simulate chat scenarios, input keys, and verify PII scrubbing."
    ]
    for tb in test_bullets:
        p_tb = tf_test.add_paragraph()
        p_tb.text = "✔ " + tb
        p_tb.font.name = "Outfit"
        p_tb.font.size = Pt(12)
        p_tb.font.color.rgb = TEXT_MUTED
        p_tb.space_before = Pt(6)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 10: Conclusion & Next Steps
    # ═══════════════════════════════════════════════════════════════
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10)
    add_header(slide10, "Summary & Next Steps")
    
    # Left Card: Summary
    create_card(slide10, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "DEMONSTRATED RESULTS")
    tb_res = slide10.shapes.add_textbox(Inches(0.95), Inches(2.3), Inches(5.3), Inches(4.0))
    tf_res = tb_res.text_frame
    tf_res.word_wrap = True
    
    p = tf_res.paragraphs[0]
    p.text = "The proof of concept successfully validates that autonomous LLM agents can be deployed in wealth portfolios with strict safety parameters:"
    p.font.name = "Outfit"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(10)
    
    res_bullets = [
        "Protects client assets via client-side scrubbing.",
        "Adheres to financial regulation via fiduciary filters.",
        "Enforces manual human gates for transfers and UTMA compliance.",
        "Tracks agent performance and latency at an admin level."
    ]
    for rb in res_bullets:
        p_rb = tf_res.add_paragraph()
        p_rb.text = "✦ " + rb
        p_rb.font.name = "Outfit"
        p_rb.font.size = Pt(12)
        p_rb.font.color.rgb = TEXT_MUTED
        p_rb.space_before = Pt(6)

    # Right Card: Next Steps
    create_card(slide10, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8), "PRODUCTION EXPANSION PATH")
    tb_ns = slide10.shapes.add_textbox(Inches(7.05), Inches(2.3), Inches(5.3), Inches(4.0))
    tf_ns = tb_ns.text_frame
    tf_ns.word_wrap = True
    
    p = tf_ns.paragraphs[0]
    p.text = "Expanding the POC beyond mock data boundaries:"
    p.font.name = "Outfit"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(10)
    
    ns_steps = [
        "Map prototype behavior to Morgan Stanley security, compliance, privacy, and model-risk requirements.",
        "Secure Entitlement Integrations: Connect to production identity profiles and access roles.",
        "Fiduciary Audit: Verify compliance of disclaimers with legal advisory groups.",
        "Pilot Deployment: Transition whitelisted sandbox connections toward limited customer testing."
    ]
    for ns in ns_steps:
        p_ns = tf_ns.add_paragraph()
        p_ns.text = "➔ " + ns
        p_ns.font.name = "Outfit"
        p_ns.font.size = Pt(11)
        p_ns.font.color.rgb = TEXT_MUTED
        p_ns.space_before = Pt(4)
        
    # Save the presentation
    prs.save("/Users/klejnieks/Graveyard/GEAP/docs/GEAP_Executive_Presentation.pptx")
    print("SUCCESS: PowerPoint saved to docs/GEAP_Executive_Presentation.pptx")

if __name__ == "__main__":
    main()
